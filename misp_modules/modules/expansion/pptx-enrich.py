import json
import binascii
import np
from pptx import Presentation
import io

misperrors = {'error': 'Error'}
mispattributes = {'input': ['attachment'],
                  'output': ['freetext', 'text']}
moduleinfo = {'version': '0.1', 'author': 'Sascha Rommelfangen',
              'description': '.pptx to freetext-import IOC extractor',
              'module-type': ['expansion']}

moduleconfig = []


def handler(q=False):
    if q is False:
        return False
    q = json.loads(q)
    filename = q['attachment']
    try:
        pptx_array = np.frombuffer(binascii.a2b_base64(q['data']), np.uint8)
    except Exception as e:
        print(e)
        err = "Couldn't fetch attachment (JSON 'data' is empty). Are you using the 'Query enrichment' action?"
        misperrors['error'] = err
        print(err)
        return misperrors

    ppt_content = ""
    ppt_file = io.BytesIO(pptx_array)
    try:
        ppt = Presentation(ppt_file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
                    ppt_content = ppt_content + "\n" + shape.text
        return {'results': [{'types': ['freetext'], 'values': ppt_content, 'comment': ".pptx-to-text from file " + filename},
                            {'types': ['text'], 'values': ppt_content, 'comment': ".pptx-to-text from file " + filename}]}
    except Exception as e:
        print(e)
        err = "Couldn't analyze file as .pptx. Error was: " + str(e)
        misperrors['error'] = err
        return misperrors


def introspection():
    return mispattributes


def version():
    moduleinfo['config'] = moduleconfig
    return moduleinfo
